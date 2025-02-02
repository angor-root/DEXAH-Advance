from pptx import Presentation
import os
import json
import time
from pptx.util import Pt

# ----------------- Configuración -----------------
texto_mayuscula = True # True para mayúscula, False para minúscula (default: False)
font_size = 24 # Tamaño de fuente (default: 24)
# ----------------- Configuración -----------------


# set yout template here, slide[0] for textbox(title), textbox(slide number), slide[1] for content
hymns = []
hymnBooks = []
# template = "a.pptx"
def main():
    print("""
░       ░░░        ░░  ░░░░  ░░░      ░░
▒  ▒▒▒▒  ▒▒  ▒▒▒▒▒▒▒▒▒  ▒▒  ▒▒▒  ▒▒▒▒  ▒
▓  ▓▓▓▓  ▓▓      ▓▓▓▓▓▓    ▓▓▓▓  ▓▓▓▓  ▓
█  ████  ██  █████████  ██  ███        █
█       ███        ██  ████  ██  ████  █
Welcome to DEXA Sophos by angor.root                                  
""")
    files = os.listdir("DEXHymnBooks")
    print("[ ] Available hymn books:")
    i = 1
    for file in files:
        with open('DEXHymnBooks/'+file, 'r', encoding='utf-8') as file:
            data = json.load(file)
            hymnBooks.append(data)
            print(f"     [{i}] {data['title']} [{data['tag']}]")
            i += 1
    print("[ ] Enter in the format 'Hymn HymnBook' e.g. '1hb'")
    print("     [!] Press 'Enter' to finish entering hymns.")
    hymns = input("Enter hymns (separated by spaces): ").split()
    return hymns
hymns = main()

final_PPTX = Presentation("base.pptx")
start_time = time.time()

for hymn in hymns:
    hymn_number = int(''.join(filter(str.isdigit, hymn)))
    hymn_book = ''.join(filter(str.isalpha, hymn))
    try:
        hymn_book = next(hymnBook for hymnBook in hymnBooks if hymnBook['tag'] == hymn_book)
        hymn_book_title = hymn_book['title']
        # Cambio en la búsqueda del himno
        himno = next((h for h in hymn_book['hymns'] if h['number'] == hymn_number), None)
        
        if himno is None:
            raise Exception(f"Himno {hymn_number} no encontrado")
            
        hymn_title = himno['title']
        print(f"[ ] Hymn book {hymn_book_title} Geting hymn {hymn_number}...")
        # ----------------- extraer coro -----------------
        coro = himno['chorus']
        # ----------------- Iterar para cada estrofa -----------------
        estrofas = himno['verses']
        capital = texto_mayuscula # True for capital, False for lowercase
        if capital:
            estrofas = [estrofa.upper() for estrofa in estrofas]
            coro = coro.upper()
            hymn_title = hymn_title.upper()            
        # ----------------- Iterar para cada himno -----------------
        # ----------------- Crear pptx -----------------
        slide_layout = final_PPTX.slide_layouts[0]  # Select the layout for a title slide
        slide = final_PPTX.slides.add_slide(slide_layout)  # Create a new slide using the selected layout

        title = slide.shapes.title  # Get the title shape of the slide
        title.text = hymn_title # Set the title text

        subtitle = slide.placeholders[1]  # Get the subtitle shape of the slide
        # subtitle.text = f"Hymn {hymn_number} from {hymn_book_title}" # ------------ [cAMBIO]
        subtitle.text = f"{hymn_number} - {hymn_book_title}"
        # ----------------- estrofas -----------------
        t = 5 # stable
        # t = 0 # test
        for estrofa in estrofas:
            slide_layout = final_PPTX.slide_layouts[t]  # Select the layout for a title slide
            slide = final_PPTX.slides.add_slide(slide_layout)  # Create a new slide using the selected layout
            title = slide.shapes.title  # Get the title shape of the slide
            title.text = estrofa  # Set the title text
            # ----------------- coro -----------------
            if coro != "":
                slide_layout = final_PPTX.slide_layouts[t]
                slide = final_PPTX.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = coro
        # ----------------- Guardar pptx -----------------
    except StopIteration:
        print(f"[ ] Hymn book {hymn_book} not found.")
        continue
    except Exception as e:
        print(f"[ ] An error occurred: {e}")
        continue

    final_PPTX.save("DEXAH_output.pptx")

print("[ ] Done! Check DEXAH_output.pptx")

# tiempo de ejecución
print("--- %s seconds ---" % (time.time() - start_time))