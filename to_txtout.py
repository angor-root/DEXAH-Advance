from pptx import Presentation
import os
import json
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
import time

def duplicate_slide(pres, slide):
    slide_id = pres.slides._sldIdLst[-1].rId
    slide = pres.slides[-1]._element
    new_slide = parse_xml(slide.toxml())
    slide_id = pres.slides._sldIdLst.get_or_add_sldId()
    slide_id.rId = pres.part.relate_to(new_slide, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
    return pres.slides[-1]

# set yout template here, slide[0] for textbox(title), textbox(slide number), slide[1] for content
hymns = []
hymnBooks = []
# template = "a.pptx"
def main():
    print("""
░       ░░░        ░░  ░░░░  ░░░      ░░
▒  ▒▒▒▒  ▒▒  ▒▒▒▒▒▒▒▒▒  ▒▒  ▒▒▒  ▒▒▒▒  ▒
▓  ▓▓▓▓  ▓▓      ▓▓▓▓▓▓    ▓▓▓▓  ▓▓▓▓  ▓
█  ████  ██  █████████  ██  ███        █
█       ███        ██  ████  ██  ████  █
Welcome to DEXA Sophos by angor.root                                  
""")
    files = os.listdir("DEXHymnBooks")
    print("[ ] Available hymn books:")
    for file in files:
        with open('DEXHymnBooks/'+file, 'r', encoding='utf-8') as file:
            data = json.load(file)
            hymnBooks.append(data)
            print(f"     [!] {data['title']} [{data['tag']}]")
    print("[ ] Enter in the format 'Hymn HymnBook' e.g. '1hb'")
    print("     [!] Press 'Enter' to finish entering hymns.")
    while True:
        i = 1
        hymn = input(f'[{i}] Hymn HymnBook: ')
        if hymn == '':
            break
        hymns.append(hymn)
        i += 1
    
main()

final_PPTX = Presentation("base.pptx")
start_time = time.time()
for hymn in hymns:
    hymn_number = int(''.join(filter(str.isdigit, hymn)))
    hymn_book = ''.join(filter(str.isalpha, hymn))
    try:
        hymn_book = next(hymnBook for hymnBook in hymnBooks if hymnBook['tag'] == hymn_book)
        hymn_book_title = hymn_book['title']
        hymn_title = hymn_book['hymns'][hymn_number-1]['title']
        # print(f"[ ] Hymn book {hymn_book_title} Geting hymn {hymn_number}...")
        # ----------------- extraer coro -----------------
        coro = hymn_book['hymns'][hymn_number-1]['chorus']
        # ----------------- Iterar para cada estrofa -----------------
        estrofas = hymn_book['hymns'][hymn_number-1]['verses']
        # ----------------- Iterar para cada himno -----------------
        # ----------------- Crear pptx -----------------
        slide_layout = final_PPTX.slide_layouts[0]  # Select the layout for a title slide
        slide = final_PPTX.slides.add_slide(slide_layout)  # Create a new slide using the selected layout

        title = slide.shapes.title  # Get the title shape of the slide
        title.text = hymn_title # Set the title text

        subtitle = slide.placeholders[1]  # Get the subtitle shape of the slide
        # subtitle.text = f"Hymn {hymn_number} from {hymn_book_title}" # ------------ [cAMBIO]
        subtitle.text = f"{hymn_number} - {hymn_book_title}"
        # ----------------- estrofas -----------------
        print("\\" + "textbf{" + f'{hymn_title} - {hymn_book_title} {hymn_number}' +"}\n")
        t = 5
        print(f"{1}. {estrofas[0]}\n")
        if coro != "":
            print(f"Coro: {coro}\n")
        for i in range(1, len(estrofas)):
            print(f"{i+1}. {estrofas[i]}\n")
        # ----------------- Guardar pptx -----------------
    except StopIteration:
        print(f"[ ] Hymn book {hymn_book} not found.")
        continue
    except Exception as e:
        print(f"[ ] An error occurred: {e}")
        continue



# tiempo de ejecución
print("--- %s seconds ---" % (time.time() - start_time))