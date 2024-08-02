from pptx import Presentation

filename = 'exact_n/hb5.pptx'  # Reemplaza esto con la ruta a tu archivo
prs = Presentation(filename)

my_list = []
slide_numbers = [0,1, 2, 3]  # List of slide numbers to evaluate

for slide_number in slide_numbers:
    slide = prs.slides[slide_number]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.hyperlink is not None:
                    index_ref = str(run.hyperlink.address)
                    number = ''.join(filter(str.isdigit, index_ref))
                    my_list.append((run.text, number))

no_number = []
i = 0
for i in range(len(my_list)):
    if my_list[i][1]:  # Comprobar si la cadena no está vacía
        i = int(my_list[i][1])
        slide_number = i - 1  # Adjusting slide number to zero-based index
        slide = prs.slides[slide_number]
        for shape in slide.shapes:
            if shape.has_text_frame:
                # si el texto no tiene número
                i += 1
                if not any(char.isdigit() for char in shape.text_frame.text):

                    print(shape.text_frame.text)
                    print('----------------')
print(i)
# abrir ./content.txt y cuardar en una lsita las 10 primeras lineas de cada texto si contiene el numero 1. :
# contenidos = []
# with open('exact_n/content.txt', 'r') as file:
#     for line in file:
#         if '1' in line:
#             contenidos.append(line)

# def comparar_content(contenidos, buscar):
#     for content in contenidos:
#         if not buscar in content:
#             print(content)